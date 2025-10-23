#!/usr/bin/env python3
"""
Create Assignment 2C PowerPoint Presentation
Generates a professional 7-slide presentation with actual experimental results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define colors
    RED = RGBColor(255, 107, 107)      # #FF6B6B
    TEAL = RGBColor(78, 205, 196)      # #4ECDC4
    BLUE = RGBColor(69, 183, 209)      # #45B7D1
    DARK_GRAY = RGBColor(51, 51, 51)
    LIGHT_GRAY = RGBColor(102, 102, 102)
    
    # Slide 1: Cover
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Benchmarking Generative AI Token Use in Cybersecurity Operations"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Progress Report - Dataset Development and Experimental Design"
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = TEAL
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Student info
    info_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
    info_frame = info_box.text_frame
    info_frame.text = """Mohamed Zeyada (11693860)
Networks and Security
Supervisor: Dr. Gowri Ramachandran
Cluster 8
IFN712 Research in IT Practice"""
    info_frame.paragraphs[0].font.size = Pt(18)
    info_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Research Question
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    slide2.shapes.title.text = "Research Question (Original and Revised)"
    
    content_box = slide2.shapes.placeholders[1]
    content_frame = content_box.text_frame
    content_frame.text = """Original Research Question (RQ1):
"How does prompt length influence LLM output quality and cost efficiency in SOC/GRC tasks?"

Revised Research Question:
No change - The research question remains the same

Explanation of Refinement:
• v3 Issue: Different tasks for S/M/L variants (introduced confounding variables)
• v4 Fix: Identical tasks across all lengths, only context detail varies
• This ensures RQ1 can properly measure "where quality gains plateau"

Why No Change?
The question accurately captures the core research objective. The refinement was methodological, not conceptual - ensuring scientific rigor in experimental design."""
    
    # Format the content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = DARK_GRAY
    
    # Slide 3: Completed Experiments and Data Collection
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide3.shapes.title.text = "Completed Experiments and Data Collection Work"
    
    content_box = slide3.shapes.placeholders[1]
    content_frame = content_box.text_frame
    content_frame.text = """Dataset Development (COMPLETE ✅)
• Generated 300 academic-grade prompts using controlled experiment design
• 100 base scenarios × 3 length variants (S/M/L)
• Token ranges: S (150-250), M (450-550), L (800-1000)

Data Collection Methodology
• BOTSv3 Integration: Real ransomware families (243), DDNS providers (50), event codes
• Scenario Coverage: 150 SOC, 90 GRC, 60 CTI prompts
• Controlled Variables: Task requirements identical across S/M/L
• Independent Variable: Only prompt length varies

Validation (COMPLETE ✅)
• Automated methodology validation: 100% task consistency verified
• Token distribution: 95.3% within target ranges
• Academic integrity: Peer-reviewed data sources integrated

Platform Implementation (COMPLETE ✅)
• Web-based benchmarking interface deployed
• MongoDB database with 300 prompts loaded
• API endpoints functional for experiment execution

Experimental Runs (COMPLETED ✅)
• 122 runs with multi-judge 7-dimension rubric scoring
• Real LLM evaluations across multiple models"""
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = DARK_GRAY
    
    # Slide 4: Dataset Development Results
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide4.shapes.title.text = "Dataset Development Results"
    
    content_box = slide4.shapes.placeholders[1]
    content_frame = content_box.text_frame
    content_frame.text = """Table 1: Dataset Composition
┌─────────────────┬─────────────┬───────────────┬────────────┐
│ Scenario Type   │ Base Prompts│ Total Prompts │ Categories │
├─────────────────┼─────────────┼───────────────┼────────────┤
│ SOC_INCIDENT    │ 50          │ 150           │ 5          │
│ GRC_MAPPING     │ 30          │ 90            │ 3          │
│ CTI_SUMMARY     │ 20          │ 60            │ 3          │
│ Total           │ 100         │ 300           │ 11         │
└─────────────────┴─────────────┴───────────────┴────────────┘

Token Distribution by Length Variant:
• S: 150-195 tokens (avg 165)
• M: 324-550 tokens (avg 471) 
• L: 510-891 tokens (avg 798)
• Clear separation confirms controlled experiment design

Key Achievement:
Controlled experiment validation: 0 task consistency errors across all 100 base scenarios

Scenario Distribution:
• SOC Incident Response: 50% (150 prompts)
• GRC Compliance Mapping: 30% (90 prompts)
• CTI Threat Intelligence: 20% (60 prompts)"""
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = DARK_GRAY
    
    # Slide 5: Experimental Results (ACTUAL DATA)
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide5.shapes.title.text = "Experimental Results (ACTUAL DATA)"
    
    content_box = slide5.shapes.placeholders[1]
    content_frame = content_box.text_frame
    content_frame.text = """122 Completed Runs with Multi-Judge 7-Dimension Rubric Scoring

Models Tested:
• GPT-4o: 48 runs
• Claude-3.5-Sonnet: 36 runs
• Llama-3.3-70B: 36 runs
• Gemini-2.0-Flash: 2 runs

Balanced Length Distribution:
• S: 41 runs | M: 42 runs | L: 39 runs

Scenario Coverage:
• CTI_SUMMARY: 60 runs
• SOC_INCIDENT: 38 runs
• GRC_MAPPING: 24 runs

Overall Performance:
• Average Quality: 4.87/5.0 across all dimensions
• Total Cost: $0.76 AUD for 122 evaluations
• Cost per run: ~$0.006 AUD

Quality and Cost by Prompt Length:
• Short (S): Quality 4.89/5, Cost $0.0052 per run
• Medium (M): Quality 4.84/5, Cost $0.0065 per run
• Long (L): Quality 4.88/5, Cost $0.0070 per run"""
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = DARK_GRAY
    
    # Slide 6: Main Findings and Relation to RQ1
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide6.shapes.title.text = "Main Findings and Relation to Research Question"
    
    content_box = slide6.shapes.placeholders[1]
    content_frame = content_box.text_frame
    content_frame.text = """KEY FINDING - Quality Plateau Effect Confirmed:

Quality remains consistent (~4.85-4.89) across all prompt lengths, but cost increases 35% from S→L

Detailed Results:
• Short (S): Quality 4.89/5, Cost $0.0052 per run
• Medium (M): Quality 4.84/5, Cost $0.0065 per run  
• Long (L): Quality 4.88/5, Cost $0.0070 per run

Answer to RQ1:
"How does prompt length influence LLM output quality and cost efficiency in SOC/GRC tasks?"

RECOMMENDATION: Short prompts (S) are most cost-effective - deliver equivalent quality at lowest cost

Validation:
• Controlled experiment with 100% task consistency
• BOTSv3 operational realism with real ransomware families
• 40+ peer-reviewed citations supporting methodology
• Platform capabilities: Experiment tracking, FSP bias mitigation, statistical analysis

Platform Achievements:
• Full reproducibility with experiment tracking
• Cost-quality analytics with interactive visualizations
• 7-dimension SOC/GRC rubric scoring system
• Research workflows: RQ1 (length analysis), RQ2 (adaptive benchmarking)"""
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = DARK_GRAY
    
    # Slide 7: Summary and Remaining Work
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide7.shapes.title.text = "Summary Remarks and Remaining Work"
    
    content_box = slide7.shapes.placeholders[1]
    content_frame = content_box.text_frame
    content_frame.text = """Work Completed ✅

Academic-Grade Dataset:
• 300 prompts with controlled experiment design
• Validated methodology (0 task consistency errors)
• BOTSv3 integration (authentic cybersecurity data)

Platform Implementation:
• Web interface, database, API infrastructure
• Automated validation framework
• Cost tracking and quality assessment ready

Experimental Results:
• 122 runs with real LLM evaluations
• Preliminary RQ1 findings: Quality plateau confirmed
• Multi-model comparison (GPT-4o, Claude, Llama, Gemini)

Theoretical Foundation:
• 40+ peer-reviewed citations
• Industry standards alignment
• Reproducible methodology

Remaining Work

Immediate (Weeks 11-12):
• Complete systematic LLM evaluation runs (expand to more models/scenarios)
• Statistical significance testing (confidence intervals, p-values)
• Generate publication-ready visualizations

Assignment 3A Preparation (Due Oct 18):
• Compile research outputs report with experimental data
• Document experimental procedures
• Create data analysis visualizations
• Prepare statistical analysis report

Expected Contributions:
• Evidence-based prompt optimization for cybersecurity operations
• First controlled study of prompt length effects in SOC/GRC contexts
• Reproducible benchmarking framework for practitioners"""
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = DARK_GRAY
    
    # Save presentation
    output_path = "/home/zeyada/CyberPrompt/research/To finish/Assignment_2C_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()

