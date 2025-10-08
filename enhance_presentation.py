#!/usr/bin/env python3
"""
Enhance Assignment 2C PowerPoint Presentation with Charts
Adds the existing token and scenario distribution charts to the presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def enhance_presentation():
    # Load existing presentation
    prs = Presentation("/home/zeyada/CyberPrompt/research/To finish/Assignment_2C_Presentation.pptx")
    
    # Define colors
    RED = RGBColor(255, 107, 107)      # #FF6B6B
    TEAL = RGBColor(78, 205, 196)      # #4ECDC4
    BLUE = RGBColor(69, 183, 209)      # #45B7D1
    DARK_GRAY = RGBColor(51, 51, 51)
    
    # Add charts to Slide 4 (Dataset Development Results)
    slide4 = prs.slides[3]  # 0-indexed
    
    # Add token distribution chart
    token_chart_path = "/home/zeyada/CyberPrompt/research/To finish/token_distribution_chart.png"
    if os.path.exists(token_chart_path):
        # Add chart on the right side
        slide4.shapes.add_picture(token_chart_path, Inches(5.5), Inches(2), Inches(4), Inches(3))
        
        # Add chart title
        chart_title = slide4.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(0.5))
        chart_title.text_frame.text = "Token Distribution by Length Variant"
        chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart_title.text_frame.paragraphs[0].font.bold = True
        chart_title.text_frame.paragraphs[0].font.color.rgb = TEAL
        chart_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add scenario distribution chart
    scenario_chart_path = "/home/zeyada/CyberPrompt/research/To finish/scenario_distribution_chart.png"
    if os.path.exists(scenario_chart_path):
        # Add chart below the first one
        slide4.shapes.add_picture(scenario_chart_path, Inches(5.5), Inches(5.5), Inches(4), Inches(3))
        
        # Add chart title
        chart_title2 = slide4.shapes.add_textbox(Inches(5.5), Inches(5), Inches(4), Inches(0.5))
        chart_title2.text_frame.text = "Scenario Distribution"
        chart_title2.text_frame.paragraphs[0].font.size = Pt(14)
        chart_title2.text_frame.paragraphs[0].font.bold = True
        chart_title2.text_frame.paragraphs[0].font.color.rgb = TEAL
        chart_title2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Modify Slide 4 content to fit with charts
    # Remove the existing content box and add a smaller one on the left
    for shape in slide4.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame.text:
            if "Table 1: Dataset Composition" in shape.text_frame.text:
                # Resize and reposition the content box
                shape.left = Inches(0.5)
                shape.top = Inches(1.5)
                shape.width = Inches(4.5)
                shape.height = Inches(7)
                break
    
    # Add a summary box for Slide 5 with key findings
    slide5 = prs.slides[4]  # Experimental Results slide
    
    # Add a highlighted findings box
    findings_box = slide5.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1.5))
    findings_frame = findings_box.text_frame
    findings_frame.text = """KEY FINDING: Quality Plateau Effect Confirmed
Quality remains consistent (4.84-4.89/5) across all prompt lengths, but cost increases 35% from Short→Long
RECOMMENDATION: Use Short prompts for optimal cost-effectiveness in SOC/GRC operations"""
    
    # Format the findings box
    for paragraph in findings_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RED
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add border to findings box
    findings_box.line.color.rgb = RED
    findings_box.line.width = Pt(2)
    
    # Save enhanced presentation
    output_path = "/home/zeyada/CyberPrompt/research/To finish/Assignment_2C_Presentation_Enhanced.pptx"
    prs.save(output_path)
    print(f"Enhanced presentation saved to: {output_path}")
    
    return output_path

if __name__ == "__main__":
    enhance_presentation()

